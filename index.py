# /home/dunghm/Du_an_sinh_video_main_goloi/index.py
import os
import tempfile
import logging
import gradio as gr
from shutil import which
import asyncio
from gtts import gTTS
try:
    import edge_tts
except Exception:
    edge_tts = None
# Giữ theo môi trường của bạn
LIBREOFFICE_APPIMAGE = "/home/dunghm/LibreOffice-still.basic-x86_64.AppImage"

# ==== VOICE (Edge TTS + XTTS clone) ====
from src.utils.xtts_clone import create_cloned_voice, list_supported_languages

logger = logging.getLogger(__name__)

# Bản đồ voice có sẵn (Edge TTS)
EDGE_VOICE_BY_LANG_GENDER = {
    "vi": {"Nữ": "vi-VN-HoaiMyNeural",      "Nam": "vi-VN-NamMinhNeural"},
    "en": {"Nữ": "en-US-JennyNeural",        "Nam": "en-US-GuyNeural"},
    "zh": {"Nữ": "zh-CN-XiaoxiaoNeural",     "Nam": "zh-CN-YunxiNeural"},
    "ja": {"Nữ": "ja-JP-NanamiNeural",       "Nam": "ja-JP-KeitaNeural"},
    "ko": {"Nữ": "ko-KR-SunHiNeural",        "Nam": "ko-KR-InJoonNeural"},
    "fr": {"Nữ": "fr-FR-DeniseNeural",       "Nam": "fr-FR-HenriNeural"},
    "de": {"Nữ": "de-DE-KatjaNeural",        "Nam": "de-DE-ConradNeural"},
    "es": {"Nữ": "es-ES-ElviraNeural",       "Nam": "es-ES-AlvaroNeural"},
    "it": {"Nữ": "it-IT-ElsaNeural",         "Nam": "it-IT-IsmaelNeural"},
    "pt": {"Nữ": "pt-BR-FranciscaNeural",    "Nam": "pt-BR-AntonioNeural"},
}

def convert_text_to_audio(text, language='vi', gender='Nữ', preferred_voice: str | None = None):
    """
    TTS ưu tiên Edge TTS (nếu có voice), fallback gTTS. Trả về đường dẫn file mp3 tạm.
    """
    try:
        language = language or 'vi'
        if not text or not text.strip():
            return None

        tmp = tempfile.NamedTemporaryFile(delete=False, suffix='.mp3')
        out_path = tmp.name
        tmp.close()

        # chọn voice
        voice = preferred_voice or _get_edge_voice(language, gender)
        if edge_tts is not None and voice:
            try:
                async def _save():
                    communicate = edge_tts.Communicate(text=text, voice=voice, rate="+0%", volume="+0%")
                    await communicate.save(out_path)
                asyncio.run(_save())
                return out_path
            except Exception:
                pass  # fallback

        # Fallback gTTS
        tts = gTTS(text=text, lang=language, slow=False)
        tts.save(out_path)
        return out_path
    except Exception:
        return None

def _get_edge_voice(lang_code: str, gender_label: str) -> str | None:
    try:
        return EDGE_VOICE_BY_LANG_GENDER.get(lang_code, {}).get(gender_label)
    except Exception:
        return None

def _on_voice_mode_change(mode):
    use_builtin = mode == 'Ngôn ngữ có sẵn'
    use_clone = mode == 'Giọng nhân bản'
    return gr.update(visible=use_builtin), gr.update(visible=use_clone)

def _on_builtin_lang_or_gender_change(lang, gender):
    voice = _get_edge_voice(lang, gender)
    choices = [voice] if voice else []
    return gr.update(choices=choices, value=voice if voice else None)

def _get_cloned_voice_options(root_dir='./cloned_voices'):
    options = []
    try:
        if not os.path.isdir(root_dir):
            return options
        for voice_id in sorted(os.listdir(root_dir)):
            vdir = os.path.join(root_dir, voice_id)
            if not os.path.isdir(vdir): 
                continue
            cfg = os.path.join(vdir, "config.json")
            if os.path.isfile(cfg):
                import json
                try:
                    with open(cfg, "r", encoding="utf-8") as f:
                        data = json.load(f)
                    dn = data.get("display_name")
                    if dn:
                        options.append(dn)
                        continue
                except Exception:
                    pass
            # fallback theo tên mp3 đầu tiên
            for fn in os.listdir(vdir):
                if fn.lower().endswith(".mp3"):
                    options.append(os.path.splitext(fn)[0])
                    break
    except Exception:
        pass
    return options

def _handle_create_clone(file):
    if file is None:
        return "❌ Vui lòng tải lên file mp3 giọng mẫu", gr.update()
    ok, name, err = create_cloned_voice(file.name, voices_root='./cloned_voices')
    if not ok:
        return f"❌ Tạo giọng clone thất bại: {err}", gr.update()
    opts = _get_cloned_voice_options()
    new_value = name if name in opts else None
    return f"✅ Đã tạo giọng clone: {name}", gr.update(choices=opts, value=new_value)

# ==== PPTX → ảnh + trích xuất text (tận dụng hạ tầng sẵn có) ====
from pptx import Presentation
from pdf2image import convert_from_path
from PIL import Image
import pytesseract
import zipfile, xml.etree.ElementTree as ET
from src.utils.math_formula_processor import MathFormulaProcessor, process_math_text

def _convert_pptx_to_images(pptx_path, dpi=220):
    if which(LIBREOFFICE_APPIMAGE) is None:
        raise RuntimeError("Không tìm thấy LibreOffice AppImage. Kiểm tra đường dẫn.")
    tmpdir = tempfile.mkdtemp(prefix="pptx2img_")
    import subprocess
    subprocess.run(
        [LIBREOFFICE_APPIMAGE, "--headless", "--convert-to", "pdf", "--outdir", tmpdir, pptx_path],
        check=True, stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL
    )
    pdf_path = os.path.join(tmpdir, os.path.splitext(os.path.basename(pptx_path))[0] + ".pdf")
    images = convert_from_path(pdf_path, dpi=dpi, output_folder=tmpdir, fmt='png')
    out = []
    for i, img in enumerate(images, 1):
        p = os.path.join(tmpdir, f"slide-{i:02d}.png")
        img.save(p)
        out.append(p)
    return out

# --- một bản trích xuất gọn: ưu tiên math_processor, rơi xuống đọc thô + OCR ---
def _as_path(p):
    """Trả về đường dẫn file từ path string / os.PathLike / gradio file object."""
    if p is None:
        return None
    if isinstance(p, (str, os.PathLike)):
        return str(p)
    # gr.File thường là một object có .name
    return getattr(p, "name", None)

def _extract_slides_from_pptx(pptx_file):
    """
    Chấp nhận: đường dẫn (str/PathLike) hoặc đối tượng file có .name
    Trả về: list[ {slide_number, text, image_path, has_math_objects} ]
    """
    pptx_path = _as_path(pptx_file)
    if not pptx_path or not os.path.exists(pptx_path):
        raise RuntimeError("Không tìm thấy file PowerPoint hợp lệ.")

    # convert PPTX -> ảnh PNG
    imgs = _convert_pptx_to_images(pptx_path, dpi=220)

    slides = []
    # Ưu tiên MathFormulaProcessor
    mp = MathFormulaProcessor()
    res = mp.process_powerpoint_text(pptx_path)
    if not res.get("error"):
        for s in res["slides"]:
            idx = s["slide_number"] - 1
            slides.append({
                "slide_number": s["slide_number"],
                "text": s["processed_text"],
                "image_path": imgs[idx] if 0 <= idx < len(imgs) else None,
                "has_math_objects": s["has_math_objects"]
            })
    else:
        # Fallback đọc thô + OCR
        from pptx import Presentation
        from PIL import Image
        try:
            import pytesseract
        except Exception:
            pytesseract = None

        prs = Presentation(pptx_path)
        for i, slide in enumerate(prs.slides):
            chunks = []
            for shp in slide.shapes:
                try:
                    if getattr(shp, "has_text_frame", False) and shp.text_frame:
                        t = (shp.text_frame.text or "").strip()
                        if t:
                            chunks.append(t)
                except Exception:
                    pass
            text = process_math_text("\n".join(chunks).strip())
            if not text and i < len(imgs) and pytesseract is not None:
                try:
                    ocr = pytesseract.image_to_string(Image.open(imgs[i]), lang="vie+eng")
                    text = process_math_text((ocr or "").strip())
                except Exception:
                    text = ""
            slides.append({
                "slide_number": i + 1,
                "text": text,
                "image_path": imgs[i] if i < len(imgs) else None,
                "has_math_objects": False
            })
    return slides


def _format_slides_as_text(slides):
    lines = []
    for s in slides:
        n = s.get("slide_number")
        lines.append(f"## Slide {n}")
        lines.append((s.get("text") or "").strip())
        lines.append("")
    return "\n".join(lines).strip()

# ===== UI chính của trang Index =====
def create_index_interface(app_state: gr.State):
    # 1 cột ở giữa (đã có CSS .index-center như lần trước)
    with gr.Column(variant="panel", elem_classes=["index-center"]):
        gr.Markdown("### 👨‍🏫 Ảnh Giảng Viên")
        src_image = gr.Image(label="Ảnh khuôn mặt", source="upload", type="filepath")

        gr.Markdown("### 📊 PowerPoint")
        pptx_file = gr.File(label="Chọn file .pptx", file_types=[".pptx"])

        gr.Markdown("### 🔊 Chế độ giọng đọc")
        voice_mode = gr.Radio(choices=["Ngôn ngữ có sẵn", "Giọng nhân bản"], label="Chế độ", value=None)

        # --- Built-in voices ---
        builtin_block = gr.Group(visible=False)
        with builtin_block:
            audio_lang = gr.Dropdown(
                choices=["vi","en","zh","ja","ko","fr","de","es","it","pt"],
                value="vi", label="Ngôn ngữ"
            )
            builtin_gender = gr.Radio(choices=["Nữ","Nam"], value="Nữ", label="Giọng")
            builtin_voice = gr.Dropdown(choices=[], value=None, label="Giọng cụ thể (tùy chọn)")
            audio_lang.change(_on_builtin_lang_or_gender_change, [audio_lang, builtin_gender], [builtin_voice])
            builtin_gender.change(_on_builtin_lang_or_gender_change, [audio_lang, builtin_gender], [builtin_voice])

        # --- Cloned voices (XTTS) ---
        cloned_block = gr.Accordion("🧬 Giọng nhân bản (XTTS-v2)", open=False, visible=False)
        with cloned_block:
            clone_mp3 = gr.File(label="Tải MP3 làm mẫu", file_types=[".mp3"])
            clone_btn = gr.Button("Tạo bản giọng clone")
            cloned_voice_list = gr.Dropdown(choices=_get_cloned_voice_options(), label="Chọn bản clone", value=None)
            cloned_lang = gr.Dropdown(choices=list_supported_languages(), value="vi", label="Ngôn ngữ nội dung")
            clone_status = gr.Textbox(label="Trạng thái", interactive=False)

        with gr.Row():
            proceed_btn = gr.Button("➡️ Tiếp tục chỉnh sửa (sang Editor)", variant="primary")

        status = gr.Markdown("", elem_id="index_status")

    # Events
    voice_mode.change(_on_voice_mode_change, [voice_mode], [builtin_block, cloned_block])
    clone_btn.click(_handle_create_clone, [clone_mp3], [clone_status, cloned_voice_list])

    def _proceed(
        state, img, pptx, vmode, alang, bgender, bvoice, clone_name, clang
    ):
        if pptx is None:
            return state, "❌ Chưa chọn file PowerPoint!"
        payload = dict(
            source_image=img,
            pptx_file=pptx.name,           # lưu path vào state
            voice_mode=vmode,
            audio_language=alang,
            builtin_gender=bgender,
            builtin_voice=bvoice,
            cloned_voice=clone_name,
            cloned_lang=clang,
        )
        state = {**(state or {}), **payload}
        return state, "✅ Đã lưu cấu hình. Sang Editor để trích xuất & chỉnh sửa nội dung."

    proceed_btn.click(
        _proceed,
        inputs=[app_state, src_image, pptx_file, voice_mode, audio_lang, builtin_gender, builtin_voice, cloned_voice_list, cloned_lang],
        outputs=[app_state, status]
    )

    return {
        "source_image": src_image,
        "pptx_file": pptx_file,
        "voice_mode": voice_mode,
        "audio_language": audio_lang,
        "builtin_gender": builtin_gender,
        "builtin_voice": builtin_voice,
        "cloned_voice_list": cloned_voice_list,
        "cloned_voice_language": cloned_lang,
        "proceed_btn": proceed_btn,
        "status": status,
        "app_state": app_state,
    }


