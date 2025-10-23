import os
import tempfile
import subprocess
from shutil import which
import gradio as gr

from pptx import Presentation
from pdf2image import convert_from_path
from PIL import Image

try:
    import pytesseract
except Exception:
    pytesseract = None

from src.utils.math_formula_processor import MathFormulaProcessor, process_math_text

# Đường dẫn LibreOffice theo môi trường của bạn
LIBREOFFICE_APPIMAGE = "/home/dunghm/LibreOffice-still.basic-x86_64.AppImage"

# ===== Preset nhanh & Save =====
def set_lecture_fast_mode():
    return 256, 'crop', False, 2, False, 0, 1.0

def _save_edited_slides_text(text: str):
    os.makedirs("results", exist_ok=True)
    path = os.path.join("results", "edited_slides.txt")
    with open(path, "w", encoding="utf-8") as f:
        f.write(text or "")
    return f"✅ Đã lưu nội dung ({len(text or '')} ký tự) vào: {path}"

# ===== Helper trích xuất =====
def _convert_pptx_to_images(pptx_path, dpi=220):
    if which(LIBREOFFICE_APPIMAGE) is None:
        raise RuntimeError("Không tìm thấy LibreOffice AppImage. Kiểm tra đường dẫn.")
    tmpdir = tempfile.mkdtemp(prefix="pptx2img_")
    subprocess.run(
        [LIBREOFFICE_APPIMAGE, "--headless", "--convert-to", "pdf", "--outdir", tmpdir, pptx_path],
        check=True, stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL
    )
    pdf_path = os.path.join(tmpdir, os.path.splitext(os.path.basename(pptx_path))[0] + ".pdf")
    images = convert_from_path(pdf_path, dpi=dpi, output_folder=tmpdir, fmt='png')
    out = []
    for i, img in enumerate(images, 1):
        p = os.path.join(tmpdir, f"slide-{i:02d}.png")
        img.save(p)
        out.append(p)
    return out

def extract_slides_from_pptx(pptx_file_or_path):
    """Chấp nhận path string hoặc object có .name"""
    pptx_path = pptx_file_or_path if isinstance(pptx_file_or_path, str) else getattr(pptx_file_or_path, "name", None)
    if not pptx_path or not os.path.exists(pptx_path):
        raise RuntimeError("Không tìm thấy file PowerPoint hợp lệ.")

    imgs = _convert_pptx_to_images(pptx_path, dpi=220)
    slides = []

    mp = MathFormulaProcessor()
    res = mp.process_powerpoint_text(pptx_path)
    if not res.get("error"):
        for s in res["slides"]:
            idx = s["slide_number"] - 1
            slides.append({
                "slide_number": s["slide_number"],
                "text": s["processed_text"],
                "image_path": imgs[idx] if 0 <= idx < len(imgs) else None,
                "has_math_objects": s["has_math_objects"]
            })
    else:
        prs = Presentation(pptx_path)
        for i, slide in enumerate(prs.slides):
            chunks = []
            for shp in slide.shapes:
                try:
                    if getattr(shp, "has_text_frame", False) and shp.text_frame:
                        t = (shp.text_frame.text or "").strip()
                        if t:
                            chunks.append(t)
                except Exception:
                    pass
            text = process_math_text("\n".join(chunks).strip())
            if not text and i < len(imgs) and pytesseract is not None:
                try:
                    ocr = pytesseract.image_to_string(Image.open(imgs[i]), lang="vie+eng")
                    text = process_math_text((ocr or "").strip())
                except Exception:
                    text = ""
            slides.append({
                "slide_number": i + 1,
                "text": text,
                "image_path": imgs[i] if i < len(imgs) else None,
                "has_math_objects": False
            })
    return slides

def _format_slides_as_text(slides):
    lines = []
    for s in slides:
        n = s.get("slide_number")
        lines.append(f"## Slide {n}")
        lines.append((s.get("text") or "").strip())
        lines.append("")
    return "\n".join(lines).strip()

# ===== UI Editor =====
def create_lecture_editor_interface(app_state: gr.State):
    with gr.Row().style(equal_height=True):
        # ===== LEFT COLUMN =====
        with gr.Column(variant='panel', elem_id="editor_left"):
            gr.Markdown("### 📝 Nội dung slide (chỉnh sửa tại đây)")

            # --- THAY vì gr.Textbox(...) ---
            slides_text = gr.Code(
                label="Nội dung slide",
                value="",
                language="markdown",   # hoặc "text"
                lines=28,
                interactive=False,     # ⬅️ mặc định chỉ đọc
                elem_id="slides_editor"
            )



            # Hai nút trên cùng một dòng
            with gr.Row():
                extract_from_state_btn = gr.Button(
                    "📂 Trích xuất từ PowerPoint (dùng file đã chọn ở Index)",
                    variant="secondary"
                )
                edit_save_btn = gr.Button("✏️ Sửa nội dung", variant="secondary")  # ⬅️ nút 2-trạng-thái

            # Status đặt dưới hai nút
            extract_status = gr.Markdown("")
            save_status = gr.Markdown("")
            # Trạng thái đang chỉnh sửa? False = xem, True = sửa
            is_editing = gr.State(False)

            def on_edit_save(editing: bool, text: str):
                """
                - Nếu chưa ở chế độ sửa -> bật sửa, đổi nút thành 'Lưu nội dung'
                - Nếu đang sửa       -> lưu file, tắt sửa, đổi nút lại 'Sửa nội dung'
                """
                if not editing:
                    # vào chế độ sửa
                    return (
                        gr.update(interactive=True),                     # slides_text -> editable
                        "✏️ Đang ở chế độ chỉnh sửa. Hãy sửa xong rồi bấm **Lưu nội dung**.",  # save_status
                        True,                                            # is_editing
                        gr.update(value="💾 Lưu nội dung")               # nút -> Lưu
                    )
                else:
                    # lưu và thoát sửa
                    msg = _save_edited_slides_text(text or "")
                    return (
                        gr.update(interactive=False),                    # slides_text -> readonly
                        msg,                                             # save_status
                        False,                                           # is_editing
                        gr.update(value="✏️ Sửa nội dung")              # nút -> Sửa
                    )

            # Gắn handler cho nút toggle
            edit_save_btn.click(
                on_edit_save,
                inputs=[is_editing, slides_text],
                outputs=[slides_text, save_status, is_editing, edit_save_btn],
            )

            # CSS/JS cho textbox cao và có scroll dọc/ngang
            # CSS + JS bắt chắc
            gr.HTML("""
            <style>
            /* Biến cột trái thành flex-column để con có thể fill 100% */
            #editor_left{
                display:flex; flex-direction:column; min-height:0;
            }
            /* 1 số version Gradio bọc thêm .gr-form/.gr-box -> cho phép co giãn */
            #editor_left .gr-form, 
            #editor_left .gr-box{
                display:flex; flex-direction:column; min-height:0;
            }

            /* Phần tử Code (gr.Code) và wrapper bên trong: cho phép fill */
            #slides_editor{
                flex:1 1 auto; display:flex; flex-direction:column; min-height:0;
            }
            #slides_editor .wrap,               /* wrapper của Code */
            #slides_editor .cm-editor,          /* container CodeMirror */
            #slides_editor .cm-scroller{        /* vùng cuộn chính */
                flex:1 1 auto;
                height:100% !important;
                min-height:0 !important;
            }

            /* Bật cuộn dọc/ngang trong editor */
            #slides_editor .cm-scroller{
                overflow:auto !important;
                overflow-x:auto !important;
            }

            /* Mobile: nếu muốn thấp hơn để không tràn */
            @media (max-width: 1024px){
                #slides_editor .cm-editor{ max-height:60vh !important; }
            }
            </style>
            """)


        # ===== RIGHT COLUMN =====
        with gr.Column(variant='panel'):
            with gr.Tabs():
                with gr.TabItem('⚙️ Cài đặt'):
                    pose_style = gr.Slider(minimum=0, maximum=46, step=1, label="Pose style", value=0)
                    size_of_image = gr.Radio([256, 512], value=256, label='Độ phân giải ảnh')
                    preprocess_type = gr.Radio(['crop','resize','full','extcrop','extfull'], value='crop', label='Xử lý ảnh')
                    is_still_mode = gr.Checkbox(label="Still Mode")
                    batch_size = gr.Slider(label="Batch size", step=1, maximum=10, value=1)
                    enhancer = gr.Checkbox(label="Dùng GFPGAN (chậm hơn)")
                    speech_rate = gr.Slider(minimum=0.6, maximum=1.6, step=0.05, value=1.0, label="Tốc độ đọc")
                    fast_mode_btn = gr.Button('⚡ Chế độ nhanh (256px, batch=2, không enhancer)', variant='secondary')
                    fast_mode_btn.click(
                        fn=set_lecture_fast_mode,
                        outputs=[size_of_image, preprocess_type, is_still_mode, batch_size, enhancer, pose_style, speech_rate]
                    )

            with gr.Tabs():
                gr.Markdown("### 🎬 Video Bài Giảng")
                final_video = gr.Video(label="Video kết quả", format="mp4").style(width=512)
                gr.Markdown("### 📊 Thông tin Video")
                info = gr.Textbox(label="Thông tin", lines=4, interactive=False)

            generate_btn = gr.Button('🎬 Tạo Video Bài Giảng', variant='primary')

    # ==== Handlers ====
    def _extract_and_fill(state):
        pptx_path = (state or {}).get("pptx_file")
        if not pptx_path:
            return gr.update(), "❌ Chưa chọn file PowerPoint ở trang Index!", state
        try:
            slides = extract_slides_from_pptx(pptx_path)
            text = _format_slides_as_text(slides)
            new_state = {**(state or {}), "slides_data": slides, "slides_text": text}
            return text, f"✅ Đã trích xuất {len(slides)} slide.", new_state
        except Exception as e:
            return gr.update(), f"❌ Lỗi trích xuất: {e}", state

    extract_from_state_btn.click(
        _extract_and_fill,
        inputs=[app_state],
        outputs=[slides_text, extract_status, app_state]
    )

    return {
        "slides_text": slides_text,
        "generate_btn": generate_btn,
        "pose_style": pose_style,
        "size_of_image": size_of_image,
        "preprocess_type": preprocess_type,
        "is_still_mode": is_still_mode,
        "batch_size": batch_size,
        "enhancer": enhancer,
        "speech_rate": speech_rate,
        "final_video": final_video,
        "info": info,
        "extract_from_state_btn": extract_from_state_btn,
        "extract_status": extract_status,
        "edit_save_btn": edit_save_btn,
        "save_status": save_status,
    }
